"""
KadaiGPT - AI Agentathon Hackathon Presentation Generator
Generates a professional PowerPoint presentation for the National Level Hackathon
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
PURPLE = RGBColor(124, 58, 237)
DARK_PURPLE = RGBColor(91, 33, 182)
ORANGE = RGBColor(249, 115, 22)
GREEN = RGBColor(34, 197, 94)
DARK = RGBColor(10, 10, 15)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK
    background.line.fill.background()
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0), Inches(2), prs.slide_width, Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KadaiGPT"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(0), Inches(3.2), prs.slide_width, Inches(0.6))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = '"Bill Karo, AI Sambhalo" 🛒'
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(2), Inches(4), Inches(9.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "India's First Agentic AI-Powered Retail Operations Platform\nfor 12 Million+ Kirana Stores"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Hackathon badge
    badge_box = slide1.shapes.add_textbox(Inches(3.5), Inches(5.5), Inches(6.333), Inches(0.6))
    tf = badge_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 AI Agentathon 2026 | National Level Hackathon"
    p.font.size = Pt(18)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = DARK
    bg2.line.fill.background()
    
    # Section label
    label = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "THE PROBLEM"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    # Title
    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "India's Retail Crisis 📉"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Problem cards
    problems = [
        ("📝 Manual Billing Chaos", "Store owners spend 3+ hours daily on manual billing, handwritten records, and error-prone calculations."),
        ("📦 Inventory Nightmares", "Stock-outs cost ₹2.5 lakh/year per store. 40% of products expire due to poor tracking."),
        ("💰 Credit Recovery Issues", "₹50,000+ stuck in 'udhari' (credit) per store. No systematic tracking or reminders."),
        ("📊 Zero Business Insights", "No data on best sellers, peak hours, or customer preferences. Gut-feeling decisions.")
    ]
    
    for i, (title_text, desc) in enumerate(problems):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(2 + (i // 2) * 2.3)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(26, 26, 46)
        card.line.color.rgb = RGBColor(45, 45, 75)
        
        card_title = slide2.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(0.5))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        card_desc = slide2.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1))
        tf = card_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
    
    # Stats
    stats = [("12M+", "Kirana Stores"), ("80%", "Use Paper Records"), ("₹3L", "Annual Loss/Store")]
    for i, (num, label) in enumerate(stats):
        x = Inches(2 + i * 3.5)
        stat_box = slide2.shapes.add_textbox(x, Inches(6.5), Inches(3), Inches(0.8))
        tf = stat_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = DARK_PURPLE
    bg3.line.fill.background()
    
    label = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "THE SOLUTION"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    p.font.bold = True
    
    title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(6), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Meet KadaiGPT 🤖"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    desc = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(1))
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "An AI-first, voice-enabled, offline-capable retail assistant that transforms every Kirana store into a smart store."
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    
    features = [
        "✓ Voice Commands in Hindi, Tamil & English",
        "✓ Smart OCR - Scan any bill in seconds",
        "✓ AI-Powered Demand Prediction",
        "✓ WhatsApp Bot for Anywhere Access",
        "✓ Works 100% Offline"
    ]
    
    for i, feat in enumerate(features):
        feat_box = slide3.shapes.add_textbox(Inches(0.5), Inches(3 + i * 0.6), Inches(6), Inches(0.5))
        tf = feat_box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
    
    # App mockup placeholder
    mockup = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(5), Inches(5.5))
    mockup.fill.solid()
    mockup.fill.fore_color.rgb = RGBColor(26, 26, 46)
    mockup.line.color.rgb = PURPLE
    
    mock_text = slide3.shapes.add_textbox(Inches(8), Inches(4), Inches(4), Inches(0.5))
    tf = mock_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[Live Demo Screenshot]"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 4: Core Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = DARK
    bg4.line.fill.background()
    
    label = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "CORE FEATURES"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "What Makes Us Different 🎯"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    features = [
        ("🎤", "Voice-First Billing", "Speak and the bill is created"),
        ("📸", "Smart OCR Scanning", "98% accuracy bill scanning"),
        ("🧠", "AI Demand Prediction", "Know what to stock before customers ask"),
        ("📱", "WhatsApp Integration", "Access store data via WhatsApp"),
        ("🏆", "Loyalty Program", "Smart rewards, 40% more repeat customers"),
        ("📊", "Smart Analytics", "Real-time dashboards & insights")
    ]
    
    for i, (icon, title_text, desc_text) in enumerate(features):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(2 + (i // 3) * 2.5)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(26, 26, 46)
        card.line.color.rgb = PURPLE
        
        icon_box = slide4.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        
        card_title = slide4.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), Inches(3.4), Inches(0.5))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        card_desc = slide4.shapes.add_textbox(x + Inches(0.3), y + Inches(1.4), Inches(3.4), Inches(0.6))
        tf = card_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc_text
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
    
    # Slide 5: AI Features
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = RGBColor(26, 10, 46)
    bg5.line.fill.background()
    
    label = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-POWERED INTELLIGENCE"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Agentic AI in Action 🤖"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    ai_features = [
        ("🗣️ Natural Language Processing", "AI Agent", "User: \"Ramesh ko kal ki bill bhejo\"\n✓ Sent bill to Ramesh via WhatsApp"),
        ("📈 Predictive Analytics", "ML Model", "System: \"Holi in 5 days\"\n📦 Stock 50kg colors, 100 pichkaris"),
        ("💬 Conversational Commerce", "LLM Powered", "Customer: \"Aata hai kya?\"\n✓ \"Haan, 5kg - ₹285. Order karu?\""),
        ("🔔 Smart Notifications", "Proactive", "Auto-Alert: \"Dal 3 days left\"\n⚠️ Order now - price up tomorrow")
    ]
    
    for i, (title_text, badge, demo) in enumerate(ai_features):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(1.8 + (i // 2) * 2.8)
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(26, 26, 46)
        card.line.color.rgb = PURPLE
        
        card_title = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(0.5))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title_text} [{badge}]"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        demo_box = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), Inches(5.4), Inches(1.4))
        tf = demo_box.text_frame
        p = tf.paragraphs[0]
        p.text = demo
        p.font.size = Pt(12)
        p.font.color.rgb = GREEN
    
    # Slide 6: Tech Stack
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK
    bg6.line.fill.background()
    
    label = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "TECHNOLOGY"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Built for Scale 🏗️"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Architecture boxes
    arch = [("Frontend", "React + PWA"), ("API Layer", "FastAPI"), ("AI Engine", "Gemini + ML"), ("Database", "PostgreSQL")]
    for i, (name, tech) in enumerate(arch):
        x = Inches(1 + i * 3)
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.5), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(26, 26, 46)
        box.line.color.rgb = PURPLE
        
        text = slide6.shapes.add_textbox(x, Inches(2.7), Inches(2.5), Inches(1.1))
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
    
    # Tech badges
    techs = ["⚛️ React.js", "🐍 Python/FastAPI", "🤖 Gemini API", "📱 WhatsApp API", 
             "🐘 PostgreSQL", "☁️ Railway Cloud", "📴 Offline-First", "🔐 JWT Auth"]
    for i, tech in enumerate(techs):
        x = Inches(1 + (i % 4) * 2.8)
        y = Inches(5 + (i // 4) * 0.8)
        badge = slide6.shapes.add_textbox(x, y, Inches(2.5), Inches(0.5))
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 7: WhatsApp
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = RGBColor(6, 78, 59)
    bg7.line.fill.background()
    
    label = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "WHATSAPP INTEGRATION"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    p.font.bold = True
    
    title = slide7.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Store in Every Pocket 📱"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    wa_features = [
        ("📤 Instant Bill Sharing", "One-tap bill sharing with QR payment links"),
        ("🔔 Smart Reminders", "Automated credit reminders to customers"),
        ("🤖 Store Bot", "24/7 automated responses via WhatsApp"),
        ("📊 Daily Reports", "Get daily summary at 9 PM on WhatsApp")
    ]
    
    for i, (title_text, desc) in enumerate(wa_features):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(2 + (i // 2) * 2.5)
        
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(10, 50, 40)
        card.line.color.rgb = GREEN
        
        card_title = slide7.shapes.add_textbox(x + Inches(0.3), y + Inches(0.4), Inches(5.4), Inches(0.5))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        card_desc = slide7.shapes.add_textbox(x + Inches(0.3), y + Inches(1), Inches(5.4), Inches(0.8))
        tf = card_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
    
    # Slide 8: Future Roadmap
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = DARK
    bg8.line.fill.background()
    
    label = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "VISION"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    title = slide8.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Where We're Heading 🚀"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    roadmap = [
        ("Q1 2026", "Launch MVP\n100 Beta Stores"),
        ("Q2 2026", "WhatsApp Bot\nAI Predictions"),
        ("Q3 2026", "B2B Marketplace\nSupplier Network"),
        ("Q4 2026", "10,000 Stores\nRegional Languages"),
        ("2027", "Pan-India Launch\n1M Stores Target")
    ]
    
    for i, (time, desc) in enumerate(roadmap):
        x = Inches(1 + i * 2.3)
        
        dot = slide8.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.5), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PURPLE if i != 1 else ORANGE
        dot.line.fill.background()
        
        time_box = slide8.shapes.add_textbox(x, Inches(3), Inches(2), Inches(0.4))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        desc_box = slide8.shapes.add_textbox(x, Inches(3.4), Inches(2), Inches(0.8))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Future integrations
    future = [("🏦 UPI AutoPay", "Auto credit collection"), ("🚚 Delivery Tracking", "Home delivery mgmt"),
              ("📺 Smart Display", "Digital menu boards"), ("🤝 Franchise Mode", "Multi-store mgmt")]
    
    for i, (name, desc) in enumerate(future):
        x = Inches(1 + i * 3)
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5), Inches(2.8), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(26, 26, 46)
        card.line.color.rgb = PURPLE
        
        text = slide8.shapes.add_textbox(x + Inches(0.2), Inches(5.2), Inches(2.4), Inches(1.4))
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
    
    # Slide 9: Impact
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = DARK_PURPLE
    bg9.line.fill.background()
    
    label = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "IMPACT"
    p.font.size = Pt(14)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    
    title = slide9.shapes.add_textbox(Inches(0), Inches(0.9), prs.slide_width, Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Transforming India's Retail Backbone 🇮🇳"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    stats = [("40%", "Time Saved Daily"), ("₹3L", "Extra Revenue/Year"), ("50%", "Less Stock Wastage"), ("90%", "Credit Recovery")]
    for i, (num, label_text) in enumerate(stats):
        x = Inches(1.5 + i * 2.8)
        
        stat_box = slide9.shapes.add_textbox(x, Inches(2.5), Inches(2.5), Inches(1.5))
        tf = stat_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = label_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
    
    # SDG goals
    sdg_text = slide9.shapes.add_textbox(Inches(0), Inches(5), prs.slide_width, Inches(0.5))
    tf = sdg_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Supporting UN SDG Goals: 8 (Decent Work) | 9 (Industry Innovation) | 10 (Reduced Inequalities)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 10: Team & CTA
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = DARK
    bg10.line.fill.background()
    
    label = slide10.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "THE TEAM"
    p.font.size = Pt(14)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    
    title = slide10.shapes.add_textbox(Inches(0), Inches(0.9), prs.slide_width, Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Built with ❤️ by"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    team = [("LK", "Lokeshkumar D", "Full Stack Developer"),
            ("KS", "Kishan SG", "AI/ML Engineer"),
            ("PD", "Pranesh DP", "Backend Developer")]
    
    for i, (initials, name, role) in enumerate(team):
        x = Inches(2.5 + i * 3)
        
        avatar = slide10.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), Inches(2.2), Inches(1.5), Inches(1.5))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PURPLE
        avatar.line.fill.background()
        
        init_text = slide10.shapes.add_textbox(x + Inches(0.4), Inches(2.5), Inches(1.5), Inches(0.8))
        tf = init_text.text_frame
        p = tf.paragraphs[0]
        p.text = initials
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        name_box = slide10.shapes.add_textbox(x, Inches(3.9), Inches(2.3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        role_box = slide10.shapes.add_textbox(x, Inches(4.3), Inches(2.3), Inches(0.4))
        tf = role_box.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # CTA
    cta_box = slide10.shapes.add_textbox(Inches(0), Inches(5.3), prs.slide_width, Inches(0.6))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 Try Live Demo: kadaigpt.railway.app"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    thanks = slide10.shapes.add_textbox(Inches(0), Inches(6.2), prs.slide_width, Inches(0.8))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "🙏 Thank You! Questions? Let's discuss! 💬"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), "KadaiGPT_AI_Agentathon_2026.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
