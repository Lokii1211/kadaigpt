"""
KadaiGPT - AI Agentathon Hackathon Presentation Generator
Uses the provided 8-slide template to create the final presentation
ALL CONTENT IN ENGLISH ONLY
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Paths
TEMPLATE_PATH = r"c:\Users\dell\Desktop\VA\AI agentathon Template (3).pptx"
OUTPUT_PATH = r"c:\Users\dell\Desktop\VA\VyaparAI\presentation\KadaiGPT_Hackathon_Final.pptx"

def fill_text_frame(text_frame, text, font_size=None, bold=None, color=None):
    """Fill a text frame with new text"""
    # Clear existing paragraphs
    for para in text_frame.paragraphs:
        para.clear()
    
    # Add new text
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color

def create_presentation():
    """Create the KadaiGPT presentation using the template"""
    
    print("Loading template...")
    prs = Presentation(TEMPLATE_PATH)
    
    print(f"Template has {len(prs.slides)} slides")
    
    # Content for 8 slides - ALL ENGLISH
    # Each slide has content mapped to its text shapes in order
    slide_contents = [
        # Slide 1 - Title/Cover
        [
            "KadaiGPT",
            "Smart Billing, AI Powered",
            "India's First Agentic AI-Powered Retail Operations Platform\nfor 12 Million+ Small Retail Stores",
            "AI Agentathon 2026 | National Level Hackathon",
            "Team: Lokeshkumar D, Kishan SG, Pranesh DP"
        ],
        # Slide 2 - Problem Statement
        [
            "THE PROBLEM",
            "India's Retail Crisis",
            "12M+ small retail stores face daily operational challenges",
            "Manual Billing Chaos - 3+ hours daily on handwritten records",
            "Inventory Nightmares - $3,000/year loss from stock-outs",
            "Credit Recovery Issues - $600+ stuck per store",
            "Zero Business Insights - No data for decision making",
            "80% still use paper records | $3,500 annual loss per store"
        ],
        # Slide 3 - Solution Overview
        [
            "THE SOLUTION",
            "Meet KadaiGPT",
            "AI-first, voice-enabled, offline-capable retail assistant",
            "Voice Commands - Speak naturally to create bills instantly",
            "Smart OCR - Scan any invoice with 98% accuracy",
            "AI Prediction - Know what to stock before customers ask",
            "WhatsApp Bot - Access your store from anywhere",
            "Offline Mode - Works without internet connection"
        ],
        # Slide 4 - Core Features
        [
            "CORE FEATURES",
            "What Makes Us Different",
            "Voice-First Billing - 'Add 2kg rice, 1L oil' creates bill instantly",
            "Smart OCR Scanning - Point camera at any invoice",
            "AI Demand Prediction - Festival-aware, weather-aware forecasting",
            "WhatsApp Integration - Send bills, get reports, manage store",
            "Loyalty Program - Increase repeat customers by 40%",
            "Smart Analytics - Real-time dashboards & profit insights"
        ],
        # Slide 5 - AI/ML Features
        [
            "AI-POWERED INTELLIGENCE",
            "Agentic AI in Action",
            "Natural Language Processing - Understands commands like 'Send yesterday's bill to customer John'",
            "Predictive Analytics - 'Festival in 5 days' triggers smart restocking recommendations",
            "Conversational Commerce - Customers chat with store bot for stock checks & orders",
            "Smart Notifications - Proactive alerts: 'Stock low, order now before price increase'"
        ],
        # Slide 6 - Technology & Architecture
        [
            "TECHNOLOGY STACK",
            "Built for Scale",
            "Frontend: React.js + PWA for offline-first experience",
            "Backend: FastAPI + Python for high-performance APIs",
            "AI Engine: Gemini API + Custom ML models",
            "Database: PostgreSQL + Redis for fast queries",
            "Infrastructure: Railway Cloud with auto-scaling",
            "Security: JWT Authentication + End-to-end encryption"
        ],
        # Slide 7 - Future Roadmap & WhatsApp
        [
            "ROADMAP & VISION",
            "Where We're Heading",
            "Q1 2026: MVP Launch with 100 Beta Stores",
            "Q2 2026: WhatsApp Bot & AI Predictions",
            "Q3 2026: B2B Marketplace & Supplier Network",
            "Q4 2026: 10,000 Stores & Regional Languages",
            "2027: Pan-India Launch targeting 1M Stores",
            "Future: UPI AutoPay | Delivery Tracking | Franchise Mode"
        ],
        # Slide 8 - Impact & Thank You
        [
            "IMPACT & TEAM",
            "Transforming Retail",
            "40% Time Saved | $3,500 Extra Revenue | 50% Less Wastage | 90% Credit Recovery",
            "UN SDG Goals: 8 (Decent Work) | 9 (Innovation) | 10 (Equality)",
            "Team: Lokeshkumar D (Full Stack) | Kishan SG (AI/ML) | Pranesh DP (Backend)",
            "Try Live Demo: kadaigpt.railway.app",
            "Thank You! Questions? Let's discuss!"
        ]
    ]
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(slide_contents):
            break
            
        content = slide_contents[slide_idx]
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        
        print(f"\nSlide {slide_idx + 1}: {len(text_shapes)} text shapes available, {len(content)} content items")
        
        # Fill text shapes with content
        for i, text in enumerate(content):
            if i < len(text_shapes):
                shape = text_shapes[i]
                
                # Determine styling based on position
                if i == 0:  # Usually section label
                    fill_text_frame(shape.text_frame, text, font_size=14, bold=True)
                elif i == 1:  # Usually main title
                    fill_text_frame(shape.text_frame, text, font_size=36, bold=True)
                else:  # Content
                    fill_text_frame(shape.text_frame, text, font_size=14, bold=False)
                    
                print(f"  Shape {i}: Set text '{text[:40]}...'")
    
    # Save
    print(f"\nSaving presentation to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("SUCCESS! Presentation created!")
    
    return OUTPUT_PATH

if __name__ == "__main__":
    output = create_presentation()
    print(f"\nFinal output: {output}")
